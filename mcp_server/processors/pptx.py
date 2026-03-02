"""PPTX processor — slide text, tables, speaker notes, and URL extraction."""

from __future__ import annotations

import io
import re
import logging
from typing import List

from pptx import Presentation

logger = logging.getLogger("mcp_server.processors.pptx")


class EnhancedPPTXTextExtractor:
    """Full-fidelity PPTX text extractor."""

    @staticmethod
    def extract_text_from_pptx(file_content: bytes) -> str:
        try:
            prs = Presentation(io.BytesIO(file_content))
            slides_text: List[str] = []

            for slide_num, slide in enumerate(prs.slides, 1):
                parts: List[str] = [f"\n=== SLIDE {slide_num} ==="]

                # Title
                if slide.shapes.title and slide.shapes.title.text.strip():
                    parts.append(f"TITLE: {slide.shapes.title.text.strip()}")

                # Body shapes
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                level = para.level or 0
                                prefix = "  " * level + ("• " if level > 0 else "")
                                parts.append(f"{prefix}{text}")

                    # Tables
                    if shape.has_table:
                        table = shape.table
                        parts.append("\n--- TABLE ---")
                        for row in table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            parts.append(" | ".join(cells))

                # Speaker notes
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"\nNOTES: {notes}")

                # URLs from relationships
                urls = EnhancedPPTXTextExtractor._extract_slide_urls(slide)
                if urls:
                    parts.append("\nLINKS:")
                    parts.extend(f"  {u}" for u in urls)

                slides_text.append("\n".join(parts))

            return "\n\n".join(slides_text)

        except Exception as exc:
            logger.error(f"PPTX extraction failed: {exc}")
            return ""

    @staticmethod
    def _extract_slide_urls(slide) -> List[str]:
        """Pull hyperlinks from slide relationships and text."""
        urls: List[str] = []
        try:
            for rel in slide.part.rels.values():
                if "hyperlink" in str(rel.reltype).lower():
                    if hasattr(rel, "target_ref") and rel.target_ref:
                        urls.append(rel.target_ref)
        except Exception:
            pass

        # Also find inline URLs
        slide_text = " ".join(
            shape.text for shape in slide.shapes if shape.has_text_frame
        )
        url_re = re.compile(r'https?://[^\s<>"\']+|www\.[^\s<>"\']+\.[^\s<>"\']+')
        urls.extend(url_re.findall(slide_text))

        return list(set(urls))
